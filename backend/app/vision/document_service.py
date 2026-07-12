"""Safe local extraction, summarization, and retrieval for supported documents."""

from __future__ import annotations

from collections.abc import Iterable
from io import BytesIO
import json
import re
from pathlib import Path
from uuid import uuid4

from docx import Document as WordDocument
from fastapi import HTTPException, UploadFile
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from backend.app.core.settings import AppSettings
from backend.app.domain.vision import (
    DocumentRecord,
    DocumentSummary,
    DocumentType,
    QuestionAnswer,
)

_EXTENSIONS = {
    ".pdf": DocumentType.PDF,
    ".docx": DocumentType.DOCX,
    ".xlsx": DocumentType.XLSX,
    ".pptx": DocumentType.PPTX,
    ".txt": DocumentType.TEXT,
    ".md": DocumentType.MARKDOWN,
    ".markdown": DocumentType.MARKDOWN,
}
_SENTENCE_BOUNDARY = re.compile(r"(?<=[.!?])\s+")
_WORD = re.compile(r"[\w'-]+")


class DocumentService:
    """Store uploads locally and expose format-neutral document operations."""

    def __init__(self, settings: AppSettings) -> None:
        self._settings = settings
        self._upload_directory = settings.vision_upload_directory
        self._index_path = self._upload_directory / "index.json"

    async def upload(self, upload: UploadFile) -> DocumentRecord:
        """Validate, persist, and extract an uploaded document."""
        filename = Path(upload.filename or "document").name
        document_type = _EXTENSIONS.get(Path(filename).suffix.lower())
        if document_type is None:
            raise HTTPException(status_code=415, detail="Unsupported document type.")

        content = await upload.read(self._settings.vision_max_upload_bytes + 1)
        if len(content) > self._settings.vision_max_upload_bytes:
            raise HTTPException(
                status_code=413, detail="Document exceeds the upload size limit."
            )
        if not content:
            raise HTTPException(
                status_code=422, detail="The uploaded document is empty."
            )

        record = self._extract(filename, document_type, content)
        self._upload_directory.mkdir(parents=True, exist_ok=True)
        stored_path = (
            self._upload_directory / f"{record.id}{Path(filename).suffix.lower()}"
        )
        stored_path.write_bytes(content)
        self._write_record(record)
        return record

    def get(self, document_id: str) -> DocumentRecord:
        """Return a document record or a consistent not-found response."""
        records = self._records()
        try:
            return records[document_id]
        except KeyError as error:
            raise HTTPException(
                status_code=404, detail="Document was not found."
            ) from error

    def list(self) -> list[DocumentRecord]:
        """List documents newest first without exposing filesystem paths."""
        return sorted(
            self._records().values(), key=lambda record: record.created_at, reverse=True
        )

    def summarize(self, record: DocumentRecord) -> DocumentSummary:
        """Create an extractive summary that remains available offline."""
        sentences = self._sentences(record.text)
        summary = (
            " ".join(sentences[:5]) or "This document contains no extractable text."
        )
        return DocumentSummary(
            summary=summary,
            character_count=len(record.text),
            table_count=len(record.tables),
        )

    def answer(self, record: DocumentRecord, question: str) -> QuestionAnswer:
        """Answer from the most relevant extracted sentences without cloud services."""
        keywords = {word.lower() for word in _WORD.findall(question) if len(word) > 2}
        candidates = self._sentences(record.text)
        ranked = sorted(
            candidates,
            key=lambda sentence: sum(
                word.lower() in keywords for word in _WORD.findall(sentence)
            ),
            reverse=True,
        )
        sources = [sentence for sentence in ranked[:3] if sentence]
        if (
            not sources
            or not keywords
            or all(
                not set(_WORD.findall(source.lower())) & keywords for source in sources
            )
        ):
            return QuestionAnswer(
                answer="I could not find an answer in the extracted document text.",
                sources=[],
            )
        return QuestionAnswer(answer=" ".join(sources[:2]), sources=sources)

    def _extract(
        self, filename: str, document_type: DocumentType, content: bytes
    ) -> DocumentRecord:
        record_id = uuid4().hex
        if document_type is DocumentType.PDF:
            reader = PdfReader(BytesIO(content))
            text = "\n".join(page.extract_text() or "" for page in reader.pages)
            return self._record(
                record_id,
                filename,
                document_type,
                content,
                text,
                page_count=len(reader.pages),
            )
        if document_type is DocumentType.DOCX:
            document = WordDocument(BytesIO(content))
            tables = [self._word_table(table.rows) for table in document.tables]
            text = "\n".join(paragraph.text for paragraph in document.paragraphs)
            return self._record(
                record_id, filename, document_type, content, text, tables=tables
            )
        if document_type is DocumentType.XLSX:
            workbook = load_workbook(BytesIO(content), data_only=True, read_only=True)
            sheet_names = workbook.sheetnames
            tables = [
                self._sheet_table(workbook[name].iter_rows(values_only=True))
                for name in sheet_names
            ]
            text = "\n\n".join(
                f"{name}\n" + self._table_text(table)
                for name, table in zip(sheet_names, tables, strict=True)
            )
            return self._record(
                record_id,
                filename,
                document_type,
                content,
                text,
                sheet_names=sheet_names,
                tables=tables,
            )
        if document_type is DocumentType.PPTX:
            presentation = Presentation(BytesIO(content))
            slides = [
                "\n".join(
                    shape.text for shape in slide.shapes if hasattr(shape, "text")
                )
                for slide in presentation.slides
            ]
            return self._record(
                record_id,
                filename,
                document_type,
                content,
                "\n\n".join(slides),
                page_count=len(slides),
            )
        text = content.decode("utf-8", errors="replace")
        return self._record(record_id, filename, document_type, content, text)

    def _record(
        self,
        record_id: str,
        filename: str,
        document_type: DocumentType,
        content: bytes,
        text: str,
        **kwargs: object,
    ) -> DocumentRecord:
        return DocumentRecord(
            id=record_id,
            filename=filename,
            document_type=document_type,
            size_bytes=len(content),
            text=text[: self._settings.vision_max_extract_characters],
            **kwargs,
        )

    @staticmethod
    def _word_table(rows: Iterable[object]) -> list[list[str]]:
        return [[cell.text.strip() for cell in row.cells] for row in rows]  # type: ignore[attr-defined]

    @staticmethod
    def _sheet_table(rows: Iterable[tuple[object, ...]]) -> list[list[str]]:
        return [["" if cell is None else str(cell) for cell in row] for row in rows]

    @staticmethod
    def _table_text(table: list[list[str]]) -> str:
        return "\n".join(" | ".join(row) for row in table)

    @staticmethod
    def _sentences(text: str) -> list[str]:
        return [
            sentence.strip()
            for sentence in _SENTENCE_BOUNDARY.split(re.sub(r"\s+", " ", text))
            if sentence.strip()
        ]

    def _records(self) -> dict[str, DocumentRecord]:
        if not self._index_path.exists():
            return {}
        raw_records = json.loads(self._index_path.read_text(encoding="utf-8"))
        return {raw["id"]: DocumentRecord.model_validate(raw) for raw in raw_records}

    def _write_record(self, record: DocumentRecord) -> None:
        records = self._records()
        records[record.id] = record
        self._index_path.write_text(
            json.dumps(
                [item.model_dump(mode="json") for item in records.values()], indent=2
            ),
            encoding="utf-8",
        )
